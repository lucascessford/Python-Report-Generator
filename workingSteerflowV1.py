import os
import pptx
import pandas
from pptx import Presentation
import openpyxl

current_dir = os.getcwd()

# 1. Finds excel file in the working directory #
def FindExcelFiles():
    if True:
        for item in os.listdir(current_dir):
            if item.endswith('.xlsx'):
                filepath = (current_dir+'\\'+item)
                return filepath
            #else:
                #print ('no excel file found')

FoundExcelFile = FindExcelFiles()
df = pandas.read_excel(FoundExcelFile, encoding='utf-8')

#####5. Presentation generator ###
prs = Presentation(current_dir + '\\ppttemplate.pptx')
 
def slidemaker(Cellreader):
    numofslides = len(Cellreader)
    for i in range(numofslides):
        slide1 = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide1)

        placeholderlist = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                placeholderlist.append(phf.idx)

        for j in range(len(placeholderlist)):
            place = slide.placeholders[placeholderlist[j]]
            place.text = str(Cellreader.iloc[:,j][i])
                
slidemaker(df)
prs.save('testpymaker.pptx')

exit






